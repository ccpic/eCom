from pyparsing import col
import numpy as np
import pandas as pd
from data_class import MonthlySalesAnalyzer
from chart_class import (
    COLOR_DICT,
    PlotLine,
    PlotStackedBar,
    PlotStackedBarPlus,
    PlotHeatGrid,
)
import matplotlib as mpl
import matplotlib.pyplot as plt
import datetime
from matplotlib.gridspec import GridSpec
from pptx import presentation, Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from datetime import datetime

QUERY_STR_MAP = {
    "ilevel_0 in ilevel_0": "总体",
    "主要客户=='京东'": "京东",
    "主要客户=='阿里'": "阿里",
    "主要客户=='其他'": "其他平台",
    "主要产品=='信立坦'": "信立坦",
    "主要产品=='泰嘉'": "泰嘉",
    "主要产品=='信立坦' and 主要客户=='京东'": "京东-信立坦",
    "主要产品=='信立坦' and 主要客户=='阿里'": "阿里-信立坦",
    "主要产品=='泰嘉' and 主要客户=='京东'": "京东-泰嘉",
    "主要产品=='泰嘉' and 主要客户=='阿里'": "阿里-泰嘉",
}
TABLE_RIGHT_2ITEMS = {
    "top": Cm(4.9),
    "left": Cm(26.5),
    "width": Cm(7),
    "height": Cm(10.7),
    "merge_cells": (
        ((0, 0), (1, 0)),
        ((2, 0), (3, 0)),
        ((4, 0), (5, 0)),
        ((6, 0), (7, 0)),
    ),
    "cell_color": {
        (0, 0): RGBColor(0, 0, 128),
        (2, 0): RGBColor(65, 105, 225),
        (4, 0): RGBColor(30, 144, 255),
        (6, 0): RGBColor(0, 191, 255),
    },
    "font_color": {
        (0, 0): RGBColor(255, 255, 255),
        (2, 0): RGBColor(255, 255, 255),
        (4, 0): RGBColor(255, 255, 255),
        (6, 0): RGBColor(255, 255, 255),
    },
}

TABLE_RIGHT_3ITEMS = {
    "top": Cm(4.9),
    "left": Cm(26.5),
    "width": Cm(7),
    "height": Cm(10.7),
    "merge_cells": (
        ((0, 0), (2, 0)),
        ((3, 0), (5, 0)),
        ((6, 0), (8, 0)),
        ((9, 0), (11, 0)),
    ),
    "cell_color": {
        (0, 0): RGBColor(0, 0, 128),
        (3, 0): RGBColor(65, 105, 225),
        (6, 0): RGBColor(30, 144, 255),
        (9, 0): RGBColor(0, 191, 255),
    },
    "font_color": {
        (0, 0): RGBColor(255, 255, 255),
        (3, 0): RGBColor(255, 255, 255),
        (6, 0): RGBColor(255, 255, 255),
        (9, 0): RGBColor(255, 255, 255),
    },
}

TEXTBOX_TABLE_TITLE = {
    "text": "贡献份额及同比变化汇总",
    "top": Cm(4.05),
    "left": Cm(27),
    "width": Cm(6),
    "height": Cm(0.86),
    "fontsize": Pt(12),
}


def get_contrib_table(
    sales: MonthlySalesAnalyzer,
    columns: str,
    values: str = "销售金额（元）",
    query_str: str = "ilevel_0 in ilevel_0",
) -> pd.DataFrame:
    """将不同时间段的不同项目贡献度及变化汇总为表格

    Parameters
    ----------
    sales : MonthlySalesAnalyzer
        一个月度销售分析模块的实例
    columns : str
        拆分数据的表头，即数据透视的列名
    values : str, optional
        数据透视的汇总值, by default "销售金额（元）"
    query_str: str, optional
        取数范围，pandas query语句的条件，类似SQL的WHERE, by default "ilevel_0 in ilevel_0",即不做筛选

    Returns
    -------
    pd.DataFrame
        一个含不同时间段贡献度及变化的pandas dataframe
    """

    dict_contrib = sales.get_contrib_kpi(
        index="年月", columns=columns, values=values, aggfunc=sum, query_str=query_str
    )
    df_contrib = pd.DataFrame()
    for k, v in dict_contrib.items():
        v.reset_index(inplace=True)
        v["周期"] = k
        v = v.reindex(columns=["周期", "index", "贡献份额", "同比"])
        v["贡献份额"] = v["贡献份额"].map("{:.1%}".format)
        v["同比"] = v["同比"].map("{:+.1%}".format)
        df_contrib = pd.concat([df_contrib, v])

    return df_contrib


class MonthlySalesPPT(object):
    """Powerpoint幻灯类，包含增加不同内容slide，并按参数添加图片和形状的方法"""

    def __init__(
        self, analyzer: MonthlySalesAnalyzer, template_path: str, save_path: str
    ) -> None:
        """初始化参数

        Parameters
        ----------
        analyzer : MonthlySalesAnalyzer
            一个月度销售分析模块的实例
        template_path : str
            PPT模板路径
        save_path : str
            新PPT保存路径
        """
        self.analyzer = analyzer
        self.template_path = template_path
        self.save_path = save_path
        self.prs = Presentation(template_path)

    def save(self):
        self.prs.save(self.save_path)
        print("PPT has been saved")

    def add_sep_slide(self, title: str = None, layout_style: int = 1, **kwargs):
        """ "添加新间隔页slide， 间隔页slide包含一个居中的大标题文本框，没有其他内容

        Parameters
        ----------
        title : str, optional
            slide标题, by default None
        layout_style : int, optional
            间隔页模板的index, by default 1
        """

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_style])

        title_placeholder = slide.placeholders[0]
        title_placeholder.text = title
        for para in title_placeholder.text_frame.paragraphs:
            font = para.runs[0].font
            if "font_size" in kwargs:
                font.size = Pt(kwargs["font_size"])
            else:
                font.size = Pt(36)

        print("Page%s" % str(int(self.prs.slides.index(slide)) + 1))

    def add_img_slide(
        self,
        imgs: dict,
        labels: list,
        title: str = None,
        layout_style: int = 0,
        *args,
        **kwargs,
    ):
        """ "添加新内容页， 内容页slide包含一个居中的图片

        Parameters
        ----------
        img: dict,
            包含要插入图片参数的字典
        title : str, optional
            slide标题, by default None
        layout_style : int, optional
            内容页模板的index, by default 0
        """

        kwargs = kwargs.copy()

        # 预设的位置和宽高参数
        IMAGE_TOP = Cm(3.5)
        IMAGE_HEIGHT = Cm(13.22)
        LABEL_TOP = Cm(2.9)
        LABEL_LEFT1 = Cm(25.23)
        LABEL_LEFT2 = Cm(28.11)
        LABEL_LEFT3 = Cm(30.99)
        LABEL_WIDTH = Cm(2.88)
        LABEL_HEIGHT = Cm(0.69)

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_style])
        shapes = slide.shapes

        # 标题
        title_shape = shapes.title
        title_shape.text = title

        # 右上角labels
        labels = [
            {
                "type": "shape",
                "text": labels[0],
                "top": LABEL_TOP,
                "left": LABEL_LEFT1,
                "width": LABEL_WIDTH,
                "height": LABEL_HEIGHT,
                "color": RGBColor(0, 128, 128),
            },
            {
                "type": "shape",
                "text": labels[1],
                "top": LABEL_TOP,
                "left": LABEL_LEFT2,
                "width": LABEL_WIDTH,
                "height": LABEL_HEIGHT,
                "color": RGBColor(220, 20, 60),
            },
            {
                "type": "shape",
                "text": labels[2],
                "top": LABEL_TOP,
                "left": LABEL_LEFT3,
                "width": LABEL_WIDTH,
                "height": LABEL_HEIGHT,
                "color": RGBColor(0, 0, 128),
            },
        ]

        for label in labels:
            obj_label = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                top=label.get("top", 0),
                left=label.get("left", 0),
                width=label.get("width", Inches(1)),
                height=label.get("height", Inches(1)),
            )
            obj_label.text = label["text"]
            obj_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            obj_label.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            obj_label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            fill = obj_label.fill
            fill.solid()  # 填充颜色前必须有此语句
            fill.fore_color.rgb = label.get("color", RGBColor(0, 0, 0))
            obj_label.line.fill.background()  # 去除边框

        # 插入图片
        for img in imgs:
            obj_img = slide.shapes.add_picture(
                image_file=img["image_file"],
                top=img.get("top", IMAGE_TOP),
                left=img.get("left", 0),
                width=img.get("width", self.prs.slide_width),
                height=img.get("height"),
            )
            if "left" in img is False:
                obj_img.left = (self.prs.slide_width - obj_img.width) / 2  # 默认图片居中

        # 插入表格
        tables = kwargs.get("tables")
        if tables is not None:
            for table in tables:
                table_data = table.get("data")

                shape_table = slide.shapes.add_table(
                    rows=table.get("rows", table_data.shape[0]),
                    cols=table.get("cols", table_data.shape[1]),
                    left=table.get("left", 0),
                    top=table.get("top", IMAGE_TOP),
                    width=table.get("width", self.prs.slide_width),
                    height=table.get("height", Cm(10)),
                )
                obj_table = shape_table.table

                # 写入数据文本
                for i, row in enumerate(obj_table.rows):
                    for j, cell in enumerate(row.cells):
                        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # 单元格文本居中（纵向）
                        obj_table.cell(i, j).text = table_data.iloc[i, j]
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER  # 单元格文本居中（横向）
                            paragraph.runs[0].font.size = table.get("fontsize", Pt(11))
                            if table_data.iloc[i, j][0] == "+":  # 正值着绿色
                                paragraph.runs[0].font.color.rgb = RGBColor(0, 176, 80)
                            elif table_data.iloc[i, j][0] == "-":  # 负值着红色
                                paragraph.runs[0].font.color.rgb = RGBColor(255, 0, 0)

                # 表格总样式
                tbl_pr = shape_table._element.graphic.graphicData.tbl
                tbl_pr[0][
                    -1
                ].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # 微软内置的样式id"NoStyleTableGrid"

                # 单元格颜色
                if "cell_color" in table:
                    for cell_xy, cell_color in table["cell_color"].items():
                        obj_table.cell(cell_xy[0], cell_xy[1]).fill.solid()
                        obj_table.cell(
                            cell_xy[0], cell_xy[1]
                        ).fill.fore_color.rgb = cell_color

                # 字体颜色
                if "font_color" in table:
                    for cell_xy, font_color in table["font_color"].items():
                        for paragraph in obj_table.cell(
                            cell_xy[0], cell_xy[1]
                        ).text_frame.paragraphs:
                            paragraph.runs[0].font.color.rgb = font_color

                # 处理合并单元格的事宜
                if "merge_cells" in table:
                    for merge_cell in table["merge_cells"]:
                        cell1_row = merge_cell[0][0]
                        cell1_col = merge_cell[0][1]
                        cell2_row = merge_cell[1][0]
                        cell2_col = merge_cell[1][1]
                        for i in range(cell2_row, cell1_row - 1, -1):
                            for j in range(cell2_col, cell1_col - 1, -1):
                                if i != cell1_row or j != cell1_col:
                                    obj_table.cell(
                                        i, j
                                    ).text = ""  # 被merge的单元格们只保留左上角主体单元格的文本
                        obj_table.cell(cell1_row, cell1_col).merge(
                            obj_table.cell(cell2_row, cell2_col)
                        )  # 根据指定左上和右下的单元格merge

                if "left" in table is False:
                    shape_table.left = (
                        self.prs.slide_width - shape_table.width
                    ) / 2  # 默认表格居中

        # 插入文本框
        textboxes = kwargs.get("textboxes")
        if textboxes is not None:
            for textbox in textboxes:
                obj_textbox = slide.shapes.add_textbox(
                    top=textbox.get("top", IMAGE_TOP),
                    left=textbox.get("left", 0),
                    width=textbox.get("width", self.prs.slide_width),
                    height=textbox.get("height", Cm(1)),
                )
                text_frame = obj_textbox.text_frame
                p = text_frame.paragraphs[0]
                p.text = textbox.get("text", "")
                p.font.size = textbox.get("fontsize", Pt(11))
                p.alignment = textbox.get("alignment", PP_ALIGN.CENTER)

        print("Page%s" % str(int(self.prs.slides.index(slide)) + 1))


def prepare_data(data_old: str, data_new: str, sheet_name: str) -> pd.DataFrame:
    """清洗数据

    Parameters
    ----------
    data_old : str
        2020&2021年数据的路径
    data_new : str
        ytd新数据的路径
    sheet_name: str
        数据中的工作表，有2个选项"达成率"和"分品规"，只有达成率表有指标

    Returns
    -------
    pd.DataFrame
        一个清洗过的pandas df
    """

    D_CLIENT_MAP = {
        "京东大药房泰州连锁有限公司": "京东",
        "阿里健康大药房医药连锁有限公司": "阿里",
        "北京德开医药科技有限公司": "北京德开",
        "广东健客医药有限公司": "广东健客",
        "广东亮健药业有限公司": "广东亮键",
        "北京好药师大药房连锁有限公司": "好药师大药房",
        "四川泉源堂药业有限公司": "四川泉源堂",
        "仁和药房网国华（北京）医药有限公司": "仁和集团",
        "广州方舟医药有限公司": "广东健客",
        "广东康爱多连锁药店有限公司": "广东康爱多",
        "广东瑞美药业有限公司": "广东瑞美",
        "阿里健康大药房连锁有限公司（浙江）": "阿里",
        "广州健客药业有限公司": "广东健客",
        "北京壹壹壹商业连锁有限公司": "北京壹壹壹",
        "叮当安健医药科技（北京）有限公司": "仁和集团",
        "阿里健康大药房医药连锁有限公司（浙江）": "阿里",
        "北京凯尔康大药房有限责任公司": "北京凯尔康",
        "广东康爱多数字健康科技有限公司": "广东康爱多",
        "广东圆心药业有限公司": "广东圆心",
        "叮当智慧药房(上海)有限公司": "叮当智慧药房",
        "北京佳康医药有限责任公司": "北京佳康",
        "京东医药（北京）有限公司": "京东",
        "广东美团大药房有限公司（美团自营）": "美团",
        "广东非常大药业有限公司（百度自营）": "百度",
        "江苏苏宁大药房有限公司（苏宁自营）": "苏宁",
        "苏州纳百特大药房有限公司（平安好医生）": "平安好医生",
        "阿里健康大药房医药连锁有限公司（信立泰旗舰店）": "阿里",
    }
    df_old = pd.read_excel(data_old, sheet_name=sheet_name, engine="openpyxl")
    df_ytd = pd.read_excel(data_new, sheet_name=sheet_name, engine="openpyxl")
    df = pd.concat([df_old, df_ytd])
    mask = df["事业部"] == "零售事业部"
    df.loc[mask, "事业部"] = "零售"
    df["年月"] = df["年月"].apply(lambda x: str(x)[:4] + "-" + str(x)[4:])
    df["主要产品"] = df["品名"].apply(lambda x: "其他" if x not in ["信立坦", "泰嘉"] else x)
    df["客户"] = df["客户名称"].map(D_CLIENT_MAP)
    df["主要客户"] = df["客户"].apply(lambda x: "其他" if x not in ["京东", "阿里"] else x)

    if sheet_name == "达成率":
        df.rename(columns={"总销量": "销售盒数", "指标数量": "指标盒数"}, inplace=True)
        df["指标金额（万元）"] = df["指标金额（元）"] / 10000
        df["指标盒数（万盒）"] = df["指标盒数"] / 10000
    else:
        df.rename(
            columns={"标准数量": "销售盒数", "当期销售金额": "销售金额（元）", "商品全称": "品规"}, inplace=True
        )

    df["销售金额（万元）"] = df["销售金额（元）"] / 10000
    df["销售盒数（万盒）"] = df["销售盒数"] / 10000

    return df


def plot_grid_metrics(
    sales: MonthlySalesAnalyzer,
    query_str: str,
) -> str:
    monthly_sales = sales.get_pivot(
        index="年月", values="销售金额（万元）", aggfunc=sum, query_str=query_str
    )
    monthly_volume = sales.get_pivot(
        index="年月", values="销售盒数（万盒）", aggfunc=sum, query_str=query_str
    )
    monthly_price = monthly_sales["销售金额（万元）"] / monthly_volume["销售盒数（万盒）"]
    monthly_price = monthly_price.to_frame(name="单价")

    plot_data = [
        monthly_sales.transpose(),
        monthly_volume.transpose(),
        monthly_price.transpose(),
    ]

    title = "电商渠道 %s 金额/盒数及单价月度趋势 - %s" % (
        QUERY_STR_MAP[query_str],
        sales.date.strftime("%Y-%m"),
    )
    gs = GridSpec(3, 1, wspace=0.1)
    fmt = [".0f", ".1f", ".1f"]

    f = plt.figure(
        width=16,
        height=5,
        FigureClass=PlotStackedBar,
        gs=gs,
        fmt=fmt,
        savepath=sales.savepath,
        data=plot_data,
        fontsize=10,
        style={
            "title": title,
            "hide_top_right_spines": True,
            "xlabel_rotation": 90,
            "last_xticks_only": True,
        },
    )

    return f.plot(show_label=False, show_total_label=True)


def plot_ratio_trend(
    sales: MonthlySalesAnalyzer,
    df_netsales: pd.DataFrame,
    product: str,
    query_str: str,
    unit: str = "金额",
    width: float = 15,
    height: float = 6,
):

    # 单位
    metric_sales = "销售金额（元）" if unit == "金额" else "销售盒数"

    df_contrib = pd.DataFrame()
    series_ecom = sales.get_pivot(
        index="年月", values=metric_sales, query_str=query_str, aggfunc=sum
    )[metric_sales]
    df_contrib["MAT"] = (
        series_ecom.rolling(12)
        .sum()
        .div(df_netsales[f"{product}-{unit}"].rolling(12).sum())
    )
    df_contrib["MQT"] = (
        series_ecom.rolling(3)
        .sum()
        .div(df_netsales[f"{product}-{unit}"].rolling(3).sum())
    )
    df_contrib["MON"] = series_ecom.div(df_netsales[f"{product}-{unit}"])
    df_contrib = df_contrib.iloc[12:, :]

    title = f"电商渠道在{product}整体发货中的贡献月度趋势 - {unit}"
    fmt = [".1%"]
    f = plt.figure(
        FigureClass=PlotLine,
        savepath=sales.savepath,
        data=df_contrib,
        fmt=fmt,
        width=width,
        height=height,
        style={
            "title": title,
            # "remove_yticks": True,
            "xlabel_rotation": 90,
            "ylabel": "电商渠道贡献份额",
        },
    )

    return f.plot(cols_showlabel=df_contrib.columns, endpoint_label_only=True)


def plot_sales_trend(
    sales: MonthlySalesAnalyzer,
    query_str: str,
    column: str,
    unit: str = "金额",
    plot_share: bool = False,
    show_title: bool = True,
    show_table: bool = False,
    show_target: bool = True,
    width: float = 16,
    height: float = 6,
    show_label: bool = True,
    show_total_label: bool = False,
    label_threshold: float = 0,
    *args,
    **kwargs,
) -> str:

    kwargs = kwargs.copy()

    # 单位
    if unit == "金额":
        metric_sales = "销售金额（万元）"
        metric_target = "指标金额（万元）"
    elif unit == "标准盒数":
        metric_sales = "销售盒数（万盒）"
        metric_target = "指标盒数（万盒）"

    monthly_sales = sales.get_pivot(
        index="年月",
        columns=column,
        values=metric_sales,
        aggfunc=sum,
        query_str=query_str,
    ).iloc[12:, :]

    if plot_share:
        monthly_sales = monthly_sales.apply(lambda x: x / x.sum(), axis=1)
        fmt = [".0%"]
        show_total_label = False
    else:
        fmt = [".0f"]

    if show_target:
        monthly_target = sales.get_pivot(
            index="年月",
            columns=column,
            values=metric_target,
            aggfunc=sum,
            query_str=query_str,
        )
        monthly_ach = monthly_sales[metric_sales] / monthly_target[metric_target]
        monthly_ach = monthly_ach.to_frame(name="达成率")
        data_line = monthly_ach.iloc[12:, :].transpose()
        fmt_line = [".0%"]
    else:
        data_line = None
        fmt_line = None

    title = kwargs.get(
        "title",
        "电商渠道销售及达成月度趋势 - %s - %s - %s"
        % (
            QUERY_STR_MAP[query_str],
            unit,
            sales.date.strftime("%Y-%m"),
        ),
    )

    if show_table:
        table_data = sales.get_kpi(query_str, unit=unit)
        fontsize = 10
    else:
        table_data = None
        fontsize = 12

    f = plt.figure(
        FigureClass=PlotStackedBar,
        width=width,
        height=height,
        savepath=sales.savepath,
        data=monthly_sales.transpose(),
        fmt=fmt,
        data_line=data_line,
        fmt_line=fmt_line,
        fontsize=fontsize,
        style={
            "title": title if show_title else None,
            "remove_yticks": True,
            "xlabel_rotation": 90,
        },
        table_data=table_data,
    )

    return f.plot(
        show_total_label=show_total_label,
        show_label=show_label,
        threshold=label_threshold,
    )


if __name__ == "__main__":
    # 数据
    df = prepare_data(
        data_old="raw data\现有架构看历史数据.xlsx",
        data_new="raw data\Q1零售事业部&电商达成率表-核对初版.xlsx",
        sheet_name="达成率",
    )
    df_ecom = df[df["事业部"] == "电商"]
    df_retail = df[df["事业部"] == "零售事业部"]
    sales = MonthlySalesAnalyzer(data=df, name="零售+电商销售", date_column="年月")
    sales_ecom = MonthlySalesAnalyzer(data=df_ecom, name="电商销售", date_column="年月")
    month_str = sales_ecom.date.strftime("%Y%m")

    df_netsales = pd.read_excel("netsales.xlsx", sheet_name="data", engine="openpyxl")
    df_netsales["年月"] = df_netsales["年月"].apply(lambda x: str(x)[:4] + "-" + str(x)[4:])
    df_netsales.set_index("年月", inplace=True)

    # 创建新ppt实例
    ppt = MonthlySalesPPT(
        sales_ecom, "./Reporting/template.pptx", "./Reporting/presentation.pptx"
    )

    # 母版
    slide = ppt.prs.slide_layouts[0]
    note = slide.shapes[5].text_frame.paragraphs[0]
    note.text = f"数据源: 电商发货数据202001-{month_str}, 当前客户看历史"
    font = note.runs[0].font
    font.italic = True
    font.size = Pt(8)

    # Page1 总标题页
    slide = ppt.prs.slides[0]
    title = slide.shapes[0].text_frame
    title.text = f"新专药\n电商渠道\n发货达成\n数据简报"
    for para in title.paragraphs:
        font = para.runs[0].font
        font.size = Pt(32)

    note = slide.shapes[2].text_frame
    note.text = month_str
    for para in note.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        font = para.runs[0].font
        font.size = Pt(16)
        font.color.rgb = RGBColor(255, 255, 255)

    author = slide.shapes[1].text_frame
    author.text = "市场研究\n" + datetime.today().strftime("%Y-%m-%d")
    for para in note.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        font = para.runs[0].font
        font.size = Pt(16)
        font.color.rgb = RGBColor(255, 255, 255)

    # Page2 - 电商渠道整体销售达成月度趋势及KPIs - 金额

    query_str = "ilevel_0 in ilevel_0"
    unit = "金额"
    ppt.add_img_slide(
        title=f"电商渠道整体销售达成月度趋势及KPIs - {unit}",
        imgs=[
            {
                "image_file": plot_sales_trend(
                    sales=sales_ecom,
                    query_str=query_str,
                    column=None,
                    show_table=True,
                    height=3,
                    unit=unit,
                    show_label=False,
                    show_total_label=True,
                ),
            }
        ],
        labels=["所有产品", "所有平台", unit],
    )

    # Page3 - 电商渠道在信立泰整体发货中的贡献月度趋势 - 金额
    COLOR_DICT["MAT"] = "teal"
    COLOR_DICT["MQT"] = "crimson"
    COLOR_DICT["MON"] = "navy"
    
    product = "信立泰"
    unit = "金额"
    ppt.add_img_slide(
        title=f"电商渠道在{product}整体发货中的贡献月度趋势 - {unit}",
        imgs=[
            {
                "image_file": plot_ratio_trend(
                    sales=sales_ecom,
                    df_netsales=df_netsales,
                    product=product,
                    query_str=query_str,
                    unit=unit,
                ),
            }
        ],
        labels=["所有产品", "所有平台", unit],
    )
    
    COLOR_DICT["MAT"] = "royalblue"
    COLOR_DICT["MQT"] = "dodgerblue"
    COLOR_DICT["MON"] = "deepskyblue"
    
    # Page4 - 零售vs.电商目标终端销售贡献月度趋势 - 金额
    column = "事业部"

    # 计算不同时间段贡献份额的kpi
    df_contrib = get_contrib_table(sales=sales, columns=column)

    ppt.add_img_slide(
        title=f"零售vs.电商目标终端销售贡献月度趋势 - {unit}",
        imgs=[
            {
                "image_file": plot_sales_trend(
                    sales=sales,
                    query_str=query_str,
                    show_target=False,
                    column=column,
                    width=13,
                    height=4,
                    show_total_label=True,
                    label_threshold=20,
                    title=f"零售vs.电商目标终端销售贡献月度趋势 - {unit}",
                ),
                "width": Cm(26),
            },
            {
                "image_file": plot_sales_trend(
                    sales=sales,
                    query_str=query_str,
                    show_target=False,
                    column=column,
                    width=13,
                    height=4,
                    plot_share=True,
                    label_threshold=0.01,
                    show_title=False,
                ),
                "width": Cm(26),
                "top": Cm(10.58),
            },
        ],
        labels=["所有产品", "零售+电商", unit],
        tables=[{**TABLE_RIGHT_2ITEMS, **{"data": df_contrib}}],
        textboxes=[TEXTBOX_TABLE_TITLE],
    )

    # Page5 - 电商渠道主要产品贡献月度趋势 - 金额
    column = "主要产品"

    # 计算不同时间段贡献份额的kpi
    df_contrib = get_contrib_table(sales=sales_ecom, columns=column)

    ppt.add_img_slide(
        title=f"电商渠道{column}贡献月度趋势 - {unit}",
        imgs=[
            {
                "image_file": plot_sales_trend(
                    sales=sales_ecom,
                    query_str=query_str,
                    show_target=False,
                    column=column,
                    width=13,
                    height=4,
                    show_total_label=True,
                    label_threshold=20,
                ),
                "width": Cm(26),
            },
            {
                "image_file": plot_sales_trend(
                    sales=sales_ecom,
                    query_str=query_str,
                    show_target=False,
                    column=column,
                    width=13,
                    height=4,
                    plot_share=True,
                    label_threshold=0.03,
                    show_title=False,
                ),
                "width": Cm(26),
                "top": Cm(10.58),
            },
        ],
        labels=["分产品", "所有平台", unit],
        tables=[{**TABLE_RIGHT_3ITEMS, **{"data": df_contrib}}],
        textboxes=[TEXTBOX_TABLE_TITLE],
    )

    # Page6 - 电商渠道主要产品贡献月度趋势 - 金额
    column = "主要客户"

    # 计算不同时间段贡献份额的kpi
    df_contrib = get_contrib_table(sales=sales_ecom, columns=column)

    ppt.add_img_slide(
        title=f"电商渠道{column}贡献月度趋势 - {unit}",
        imgs=[
            {
                "image_file": plot_sales_trend(
                    sales=sales_ecom,
                    query_str=query_str,
                    show_target=False,
                    column=column,
                    width=13,
                    height=4,
                    show_total_label=True,
                    label_threshold=20,
                ),
                "width": Cm(26),
            },
            {
                "image_file": plot_sales_trend(
                    sales=sales_ecom,
                    query_str=query_str,
                    show_target=False,
                    column=column,
                    width=13,
                    height=4,
                    plot_share=True,
                    label_threshold=0.03,
                    show_title=False,
                ),
                "width": Cm(26),
                "top": Cm(10.58),
            },
        ],
        labels=["所有产品", "分平台", unit],
        tables=[{**TABLE_RIGHT_3ITEMS, **{"data": df_contrib}}],
        textboxes=[TEXTBOX_TABLE_TITLE],
    )

    # Page7, 8, 9 - 京东平台销售达成月度趋势及KPIs - 金额
    for platform in ["京东", "阿里", "其他"]:
        query_str = f"主要客户=='{platform}'"
        ppt.add_img_slide(
            title=f"{platform}平台销售达成月度趋势及KPIs - {unit}",
            imgs=[
                {
                    "image_file": plot_sales_trend(
                        sales=sales_ecom,
                        query_str=query_str,
                        column=None,
                        show_table=True,
                        height=3,
                        unit=unit,
                        show_label=False,
                        show_total_label=True,
                    ),
                }
            ],
            labels=["所有产品", platform, unit],
        )

    # Page10-32 - 信立坦、泰嘉明细
    for product in ["信立坦", "泰嘉"]:
        # Page10 - 分隔页
        ppt.add_sep_slide(f"{product}\n明细")

        query_str = f"主要产品=='{product}'"

        # Page11 - 电商渠道在信立坦整体发货中的贡献月度趋势
        COLOR_DICT["MAT"] = "teal"
        COLOR_DICT["MQT"] = "crimson"
        COLOR_DICT["MON"] = "navy"
        
        ppt.add_img_slide(
            title=f"电商渠道在{product}整体发货中的贡献月度趋势",
            imgs=[
                {
                    "image_file": plot_ratio_trend(
                        sales=sales_ecom,
                        df_netsales=df_netsales,
                        product=product,
                        query_str=query_str,
                        unit="金额",
                        width=7,
                    ),
                    "top": Cm(4.3),
                    "left": Cm(0.5),
                    "width": None,
                    "height": Cm(11.7),
                },
                {
                    "image_file": plot_ratio_trend(
                        sales=sales_ecom,
                        df_netsales=df_netsales,
                        product=product,
                        query_str=query_str,
                        unit="标准盒数",
                        width=7,
                    ),
                    "top": Cm(4.3),
                    "left": Cm(16.5),
                    "width": None,
                    "height": Cm(11.7),
                },
            ],
            labels=[product, "所有平台", "金额/盒数"],
        )
        COLOR_DICT["MAT"] = "royalblue"
        COLOR_DICT["MQT"] = "dodgerblue"
        COLOR_DICT["MON"] = "deepskyblue"

        # Page12 - 目标产品零售vs.电商目标终端销售贡献月度趋势 - 金额

        unit = "金额"
        column = "事业部"

        # 计算不同时间段贡献份额的kpi
        df_contrib = get_contrib_table(sales=sales, columns=column, query_str=query_str)

        ppt.add_img_slide(
            title=f"{product}零售vs.电商目标终端销售贡献月度趋势 - {unit}",
            imgs=[
                {
                    "image_file": plot_sales_trend(
                        sales=sales,
                        query_str=query_str,
                        show_target=False,
                        column=column,
                        width=13,
                        height=4,
                        show_total_label=True,
                        label_threshold=20,
                        title=f"{product}零售vs.电商目标终端销售贡献月度趋势 - {unit}",
                    ),
                    "width": Cm(26),
                },
                {
                    "image_file": plot_sales_trend(
                        sales=sales,
                        query_str=query_str,
                        show_target=False,
                        column=column,
                        width=13,
                        height=4,
                        plot_share=True,
                        label_threshold=0.01,
                        show_title=False,
                    ),
                    "width": Cm(26),
                    "top": Cm(10.58),
                },
            ],
            labels=[product, "零售+电商", unit],
            tables=[{**TABLE_RIGHT_2ITEMS, **{"data": df_contrib}}],
            textboxes=[TEXTBOX_TABLE_TITLE],
        )

        # Page13, 14 目标产品销售达成月度趋势及KPIs - 金额, 标准盒数
        for unit in ["金额", "标准盒数"]:
            ppt.add_img_slide(
                title=f"{product}销售达成月度趋势及KPIs - {unit}",
                imgs=[
                    {
                        "image_file": plot_sales_trend(
                            sales=sales_ecom,
                            query_str=query_str,
                            column=None,
                            show_table=True,
                            height=3,
                            unit=unit,
                            show_label=False,
                            show_total_label=True,
                        ),
                    }
                ],
                labels=[product, "所有平台", unit],
            )

        # Page15 目标产品金额&盒数&单价月度趋势
        ppt.add_img_slide(
            title=f"{product}金额&盒数&单价月度趋势",
            imgs=[
                {
                    "image_file": plot_grid_metrics(
                        sales=sales_ecom, query_str=query_str
                    ),
                    "width": Cm(32),
                }
            ],
            labels=[product, "所有平台", "销量/价格"],
        )

        # Page16 目标产品主要客户贡献月度趋势 – 金额
        column = "主要客户"
        # 计算不同时间段贡献份额的kpi
        df_contrib = get_contrib_table(
            sales=sales_ecom, columns=column, query_str=query_str
        )

        ppt.add_img_slide(
            title=f"{product}{column}贡献月度趋势 - {unit}",
            imgs=[
                {
                    "image_file": plot_sales_trend(
                        sales=sales_ecom,
                        query_str=query_str,
                        show_target=False,
                        column=column,
                        width=13,
                        height=4,
                        show_total_label=True,
                        label_threshold=20 if product == "信立坦" else 5,
                    ),
                    "width": Cm(26),
                },
                {
                    "image_file": plot_sales_trend(
                        sales=sales_ecom,
                        query_str=query_str,
                        show_target=False,
                        column=column,
                        width=13,
                        height=4,
                        plot_share=True,
                        label_threshold=0.03,
                        show_title=False,
                    ),
                    "width": Cm(26),
                    "top": Cm(10.58),
                },
            ],
            labels=[product, "分平台", unit],
            tables=[{**TABLE_RIGHT_3ITEMS, **{"data": df_contrib}}],
            textboxes=[TEXTBOX_TABLE_TITLE],
        )

        # Page17-20 - 目标产品京东, 阿里平台销售达成月度趋势及KPIs - 金额, 标准盒数
        for platform in ["京东", "阿里"]:
            for unit in ["金额", "标准盒数"]:
                query_str = f"主要产品=='{product}' and 主要客户=='{platform}'"
                ppt.add_img_slide(
                    title=f"{product}{platform}平台销售达成月度趋势及KPIs - {unit}",
                    imgs=[
                        {
                            "image_file": plot_sales_trend(
                                sales=sales_ecom,
                                query_str=query_str,
                                column=None,
                                show_table=True,
                                height=3,
                                unit=unit,
                                show_label=False,
                                show_total_label=True,
                            ),
                        }
                    ],
                    labels=[product, platform, unit],
                )

    # Page32 - 泰嘉分品规销售金额贡献月度趋势
    df = prepare_data(
        data_old="raw data\现有架构看历史数据.xlsx",
        data_new="raw data\Q1零售事业部&电商达成率表-核对初版.xlsx",
        sheet_name="分品规",
    )
    mask = (
        (df["事业部"] == "电商")
        & (df["主要产品"] == "泰嘉")
        & (
            df["品规"].isin(
                [
                    "硫酸氢氯吡格雷片（泰嘉）25mg*20/瓶",
                    "硫酸氢氯吡格雷片（泰嘉）75mg*7片/板*4板/盒",
                    "硫酸氢氯吡格雷片（泰嘉）25mg*20/板*3板/盒",
                    "硫酸氢氯吡格雷片（泰嘉）75mg*7/盒",
                ]
            )
        )
    )
    df_ecom_talcom = df.loc[mask, :]
    sales_ecom_talcom = MonthlySalesAnalyzer(
        data=df_ecom_talcom, name="泰嘉电商销售", date_column="年月"
    )

    product = "泰嘉"
    column = "品规"
    unit = "金额"
    query_str = "ilevel_0 in ilevel_0"

    ppt.add_img_slide(
        title=f"电商渠道{product}{column}贡献月度趋势 - {unit}",
        imgs=[
            {
                "image_file": plot_sales_trend(
                    sales=sales_ecom_talcom,
                    query_str=query_str,
                    show_target=False,
                    column=column,
                    height=4,
                    show_total_label=True,
                    label_threshold=5,
                ),
                "width": Cm(32),
            },
            {
                "image_file": plot_sales_trend(
                    sales=sales_ecom_talcom,
                    query_str=query_str,
                    show_target=False,
                    column=column,
                    height=4,
                    plot_share=True,
                    label_threshold=0.03,
                    show_title=False,
                ),
                "width": Cm(32),
                "top": Cm(10.58),
            },
        ],
        labels=[product, "所有平台", unit],
    )

    ppt.save()
